#!/usr/bin/env python3
# /// script
# requires-python = ">=3.9"
# dependencies = [
#     "python-pptx>=0.6.21",
# ]
# ///
"""
Create SAP-branded PowerPoint presentations using the SAP Corp 2026 template.

Usage:
    uv run scripts/create_presentation.py --output OUTPUT_FILE [OPTIONS]

Options:
    --output FILE           Output file path (required)
    --title TEXT            Title for the first slide
    --subtitle TEXT         Subtitle for the first slide
    --layout LAYOUT         Layout for the first slide (default: "Cover K")
    --slides-json JSON      JSON array of slide definitions (use '-' to read from stdin)
    --slides-json-file FILE Read slide definitions from a JSON file (recommended for complex slides)
    --template PATH         Custom template path (default: uses SAP_Corp_2026.potx)
    --help                  Show this help message

Exit Codes:
    0 - Success
    1 - Invalid arguments
    2 - Template not found
    3 - Invalid layout name
    4 - JSON parse error
    5 - Write error

Examples:
    # Simple presentation with title slide
    uv run scripts/create_presentation.py --output presentation.pptx --title "My Presentation" --layout "Cover K"
    
    # Using a JSON file (RECOMMENDED for complex presentations)
    uv run scripts/create_presentation.py --output presentation.pptx --slides-json-file slides.json
    
    # Reading JSON from stdin
    cat slides.json | uv run scripts/create_presentation.py --output presentation.pptx --slides-json -
    
    # Inline JSON (only for simple cases without special characters)
    uv run scripts/create_presentation.py --output presentation.pptx --slides-json '[{"layout": "Cover A", "title": "Welcome"}]'

Note: For presentations with multiple slides or content containing special characters
      (newlines, quotes, etc.), use --slides-json-file to avoid shell escaping issues.
"""

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


# SAP Color definitions (RGB values)
SAP_COLORS = {
    "sap_blue": (0x1B, 0x90, 0xFF),
    "light_blue": (0x89, 0xD1, 0xFF),
    "mango": (0xE7, 0x65, 0x00),
    "teal": (0x04, 0x9F, 0x9A),
    "green": (0x36, 0xA4, 0x1D),
    "raspberry": (0xFA, 0x4F, 0x96),
    "pink": (0xF3, 0x1D, 0xED),
    "indigo": (0x78, 0x58, 0xFF),
    "black": (0x00, 0x00, 0x00),
    "white": (0xFF, 0xFF, 0xFF),
    "link_blue": (0x00, 0x70, 0xF2),
}

# Layout name to index mapping (0-based)
LAYOUT_MAP = {
    "Cover A": 0,
    "Cover B": 1,
    "Cover C": 2,
    "Cover D": 3,
    "Cover E": 4,
    "Cover F": 5,
    "Cover G": 6,
    "Cover H": 7,
    "Cover I": 8,
    "Cover J": 9,
    "Cover K": 10,
    "Cover L": 11,
    "Agenda A": 12,
    "Agenda B": 13,
    "Divider Page A": 14,
    "Divider Page B": 15,
    "Divider Page C": 16,
    "Divider Page D": 17,
    "Title Only": 18,
    "Title and Text": 19,
    "Title and Text: 2 Columns": 20,
    "Title and Text: 3 Columns": 21,
    "2 Columns - Text and Images": 22,
    "3 Columns - Text and Images": 23,
    "4 Columns - Text and Images": 24,
    "Title and Text with Image 1/3": 25,
    "Full Bleed Image": 26,
    "Text and Screenshot": 27,
    "Title and Content": 28,
    "Quote": 29,
    "Q & A": 30,
    "Thank You A": 31,
    "Thank You B": 32,
    "Blank": 33,
}


def get_script_dir() -> Path:
    """Get the directory containing this script."""
    return Path(__file__).parent.resolve()


def get_template_path(custom_path: str | None = None) -> Path:
    """Get the path to the SAP template file."""
    if custom_path:
        return Path(custom_path)
    
    # Default template location in assets folder
    # Use .pptx version (converted from .potx for python-pptx compatibility)
    script_dir = get_script_dir()
    template_path = script_dir.parent / "assets" / "SAP_Corp_2026.pptx"
    return template_path


def get_layout_index(layout_name: str) -> int:
    """Get the layout index from layout name."""
    if layout_name in LAYOUT_MAP:
        return LAYOUT_MAP[layout_name]
    
    # Try case-insensitive match
    layout_lower = layout_name.lower()
    for name, index in LAYOUT_MAP.items():
        if name.lower() == layout_lower:
            return index
    
    return -1


def add_slide(prs: Presentation, layout_name: str, title: str | None = None,
              subtitle: str | None = None, body: str | None = None,
              body2: str | None = None) -> None:
    """Add a slide with the specified layout and content."""
    layout_index = get_layout_index(layout_name)
    if layout_index < 0:
        raise ValueError(f"Invalid layout name: {layout_name}")

    if layout_index >= len(prs.slide_layouts):
        raise ValueError(f"Layout index {layout_index} out of range. Template has {len(prs.slide_layouts)} layouts.")

    slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)

    # Set title if provided and placeholder exists
    if title:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                # Title placeholder types: TITLE (1), CENTER_TITLE (3)
                if ph_type in (1, 3):
                    shape.text = title
                    break

    # Set subtitle if provided
    if subtitle:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                # Subtitle placeholder type: SUBTITLE (4)
                if ph_type == 4:
                    shape.text = subtitle
                    break

    # Set body text — collect all BODY/OBJECT placeholders in idx order,
    # then fill first with body, second with body2.
    body_shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            if ph_type in (2, 7):  # BODY (2), OBJECT (7)
                body_shapes.append(shape)
    # Sort by placeholder idx so column order is deterministic
    body_shapes.sort(key=lambda s: s.placeholder_format.idx)

    if body and len(body_shapes) >= 1:
        body_shapes[0].text = body
    if body2 and len(body_shapes) >= 2:
        body_shapes[1].text = body2


def create_presentation(
    output_path: str,
    template_path: Path,
    title: str | None = None,
    subtitle: str | None = None,
    layout: str = "Cover A",
    slides_json: str | None = None,
) -> dict:
    """Create a new SAP-branded presentation."""
    
    # Load template
    if not template_path.exists():
        raise FileNotFoundError(f"Template not found: {template_path}")
    
    prs = Presentation(str(template_path))
    
    # Remove any existing slides from the template
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    
    slides_added = []
    
    # If slides_json is provided, use it
    if slides_json:
        try:
            slides_data = json.loads(slides_json)
        except json.JSONDecodeError as e:
            raise ValueError(f"Invalid JSON: {e}")
        
        for slide_def in slides_data:
            slide_layout = slide_def.get("layout", "Title and Text")
            slide_title = slide_def.get("title")
            slide_subtitle = slide_def.get("subtitle")
            slide_body = slide_def.get("body")
            slide_body2 = slide_def.get("body2")

            add_slide(prs, slide_layout, slide_title, slide_subtitle, slide_body, slide_body2)
            slides_added.append({
                "layout": slide_layout,
                "title": slide_title,
            })
    
    # Otherwise, create a single slide if title is provided
    elif title:
        add_slide(prs, layout, title, subtitle)
        slides_added.append({
            "layout": layout,
            "title": title,
        })
    
    # If no content specified, create a blank presentation (no slides)
    # This is useful when you want to programmatically add slides later
    
    # Save the presentation
    prs.save(output_path)
    
    return {
        "output_file": output_path,
        "template_used": str(template_path),
        "slides_count": len(prs.slides),
        "slides_added": slides_added,
    }


def main():
    parser = argparse.ArgumentParser(
        description="Create SAP-branded PowerPoint presentations.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  %(prog)s --output presentation.pptx
  %(prog)s --output presentation.pptx --title "My Presentation" --layout "Cover A"
  %(prog)s --output presentation.pptx --slides-json '[{"layout": "Cover A", "title": "Welcome"}]'
        """,
    )
    parser.add_argument(
        "--output",
        required=True,
        help="Output file path",
    )
    parser.add_argument(
        "--title",
        help="Title for the first slide",
    )
    parser.add_argument(
        "--subtitle",
        help="Subtitle for the first slide",
    )
    parser.add_argument(
        "--layout",
        default="Cover K",
        help="Layout for the first slide (default: Cover K)",
    )
    parser.add_argument(
        "--slides-json",
        help="JSON array of slide definitions (use '-' to read from stdin)",
    )
    parser.add_argument(
        "--slides-json-file",
        help="Path to JSON file containing slide definitions (recommended for complex presentations)",
    )
    parser.add_argument(
        "--template",
        help="Custom template path",
    )
    
    args = parser.parse_args()
    
    # Get template path
    template_path = get_template_path(args.template)
    
    # Validate template exists
    if not template_path.exists():
        print(f"Error: Template not found: {template_path}", file=sys.stderr)
        print(f"       Ensure SAP_Corp_2026.potx is in the assets folder.", file=sys.stderr)
        sys.exit(2)
    
    # Validate layout name if not using slides_json
    if not args.slides_json and args.title:
        if get_layout_index(args.layout) < 0:
            print(f"Error: Invalid layout name: {args.layout}", file=sys.stderr)
            print(f"       Use 'uv run scripts/list_layouts.py' to see available layouts.", file=sys.stderr)
            sys.exit(3)
    
    # Determine slides JSON source (file, stdin, or direct argument)
    slides_json = None
    
    if args.slides_json_file:
        # Read from file
        json_file_path = Path(args.slides_json_file)
        if not json_file_path.exists():
            print(f"Error: JSON file not found: {args.slides_json_file}", file=sys.stderr)
            sys.exit(1)
        try:
            slides_json = json_file_path.read_text(encoding="utf-8")
        except Exception as e:
            print(f"Error: Failed to read JSON file: {e}", file=sys.stderr)
            sys.exit(1)
    elif args.slides_json:
        if args.slides_json == "-":
            # Read from stdin
            try:
                slides_json = sys.stdin.read()
            except Exception as e:
                print(f"Error: Failed to read from stdin: {e}", file=sys.stderr)
                sys.exit(1)
        else:
            slides_json = args.slides_json
    
    # Validate slides_json if provided
    if slides_json:
        try:
            slides_data = json.loads(slides_json)
            for slide_def in slides_data:
                layout_name = slide_def.get("layout", "Title and Text")
                if get_layout_index(layout_name) < 0:
                    print(f"Error: Invalid layout name in JSON: {layout_name}", file=sys.stderr)
                    sys.exit(3)
        except json.JSONDecodeError as e:
            print(f"Error: Invalid JSON: {e}", file=sys.stderr)
            print(file=sys.stderr)
            print("Tip: For complex presentations, save your JSON to a file and use --slides-json-file:", file=sys.stderr)
            print("     uv run scripts/create_presentation.py --output out.pptx --slides-json-file slides.json", file=sys.stderr)
            print(file=sys.stderr)
            print("Or pipe JSON via stdin:", file=sys.stderr)
            print("     cat slides.json | uv run scripts/create_presentation.py --output out.pptx --slides-json -", file=sys.stderr)
            sys.exit(4)
    
    try:
        result = create_presentation(
            output_path=args.output,
            template_path=template_path,
            title=args.title,
            subtitle=args.subtitle,
            layout=args.layout,
            slides_json=slides_json,
        )
        
        print(json.dumps(result, indent=2))
        print(f"\nPresentation created: {args.output}", file=sys.stderr)
        
    except FileNotFoundError as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(2)
    except ValueError as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(3)
    except Exception as e:
        print(f"Error: Failed to create presentation: {e}", file=sys.stderr)
        sys.exit(5)


if __name__ == "__main__":
    main()
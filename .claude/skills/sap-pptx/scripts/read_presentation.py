#!/usr/bin/env python3
# /// script
# requires-python = ">=3.9"
# dependencies = [
#     "python-pptx>=0.6.21",
# ]
# ///
"""
Read and extract information from PowerPoint presentations.

Usage:
    uv run scripts/read_presentation.py INPUT_FILE [OPTIONS]

Options:
    --format FORMAT         Output format: text, json, or summary (default: summary)
    --text-only             Extract only text content (no metadata)
    --slide NUM             Get info for a specific slide (1-based)
    --include-notes         Include speaker notes in output (default: true)
    --output FILE           Write output to file instead of stdout
    --help                  Show this help message

Exit Codes:
    0 - Success
    1 - Invalid arguments
    2 - File not found
    3 - Read error

Examples:
    # Get presentation summary
    uv run scripts/read_presentation.py presentation.pptx

    # Output as JSON (includes all slide data)
    uv run scripts/read_presentation.py presentation.pptx --format json

    # Extract text content only
    uv run scripts/read_presentation.py presentation.pptx --text-only

    # Get specific slide info
    uv run scripts/read_presentation.py presentation.pptx --slide 2

    # Save output to file
    uv run scripts/read_presentation.py presentation.pptx --format json --output slides.json
"""

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def extract_slide_content(slide, slide_number: int) -> dict:
    """Extract all content from a single slide including text, notes, and metadata."""
    slide_info = {
        "slide_number": slide_number,
        "title": "",
        "subtitle": "",
        "content": [],
        "speaker_notes": "",
        "shapes_count": len(slide.shapes),
        "has_notes": False,
        "layout_name": slide.slide_layout.name if slide.slide_layout else "Unknown",
    }
    
    # Extract title
    if slide.shapes.title:
        slide_info["title"] = slide.shapes.title.text.strip()
    
    # Extract text content from all shapes
    for shape in slide.shapes:
        # Skip the title shape (already captured)
        if shape == slide.shapes.title:
            continue
            
        if hasattr(shape, "text") and shape.text.strip():
            text = shape.text.strip()
            
            # Try to identify subtitle (usually the second text shape on cover slides)
            if shape.is_placeholder:
                try:
                    ph_type = shape.placeholder_format.type
                    # Subtitle placeholder type: SUBTITLE (4)
                    if ph_type == 4:
                        slide_info["subtitle"] = text
                        continue
                except:
                    pass
            
            slide_info["content"].append(text)
        
        # Handle tables
        if shape.has_table:
            table_data = []
            for row in shape.table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            if table_data:
                slide_info["content"].append({
                    "type": "table",
                    "data": table_data
                })
        
        # Handle grouped shapes
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for grouped_shape in shape.shapes:
                if hasattr(grouped_shape, "text") and grouped_shape.text.strip():
                    slide_info["content"].append(grouped_shape.text.strip())
    
    # Extract speaker notes
    if slide.has_notes_slide:
        notes_slide = slide.notes_slide
        if notes_slide.notes_text_frame:
            notes_text = notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_info["speaker_notes"] = notes_text
                slide_info["has_notes"] = True
    
    return slide_info


def extract_presentation_content(pptx_path: str, include_notes: bool = True) -> dict:
    """Extract all content from a PPTX file."""
    path = Path(pptx_path)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {pptx_path}")
    
    prs = Presentation(pptx_path)
    
    # Get presentation metadata
    core_properties = prs.core_properties
    
    presentation_info = {
        "file_path": str(path.absolute()),
        "file_name": path.name,
        "total_slides": len(prs.slides),
        "slide_width_inches": prs.slide_width.inches if prs.slide_width else None,
        "slide_height_inches": prs.slide_height.inches if prs.slide_height else None,
        "metadata": {
            "title": core_properties.title or "",
            "author": core_properties.author or "",
            "subject": core_properties.subject or "",
            "created": str(core_properties.created) if core_properties.created else "",
            "modified": str(core_properties.modified) if core_properties.modified else "",
        },
        "slides": [],
    }
    
    # Extract content from each slide
    for i, slide in enumerate(prs.slides, 1):
        slide_info = extract_slide_content(slide, i)
        
        # Remove speaker notes if not requested
        if not include_notes:
            slide_info.pop("speaker_notes", None)
            slide_info.pop("has_notes", None)
        
        presentation_info["slides"].append(slide_info)
    
    # Summary statistics
    slides_with_notes = sum(1 for s in presentation_info["slides"] if s.get("has_notes", False))
    presentation_info["summary"] = {
        "total_slides": len(prs.slides),
        "slides_with_notes": slides_with_notes if include_notes else None,
        "slides_with_titles": sum(1 for s in presentation_info["slides"] if s["title"]),
    }
    
    return presentation_info


def format_as_text(presentation_info: dict, text_only: bool = False) -> str:
    """Format presentation info as human-readable text."""
    lines = []
    
    if not text_only:
        lines.append(f"File: {presentation_info['file_name']}")
        lines.append(f"Total Slides: {presentation_info['total_slides']}")
        if presentation_info['metadata']['title']:
            lines.append(f"Title: {presentation_info['metadata']['title']}")
        if presentation_info['metadata']['author']:
            lines.append(f"Author: {presentation_info['metadata']['author']}")
        lines.append("=" * 60)
        lines.append("")
    
    for slide in presentation_info['slides']:
        if not text_only:
            lines.append(f"--- Slide {slide['slide_number']} ({slide['layout_name']}) ---")
        
        if slide['title']:
            if text_only:
                lines.append(slide['title'])
            else:
                lines.append(f"Title: {slide['title']}")
        
        if slide.get('subtitle'):
            if text_only:
                lines.append(slide['subtitle'])
            else:
                lines.append(f"Subtitle: {slide['subtitle']}")
        
        if slide['content']:
            if not text_only:
                lines.append("Content:")
            for item in slide['content']:
                if isinstance(item, dict) and item.get('type') == 'table':
                    for row in item['data']:
                        lines.append("  " + " | ".join(row))
                else:
                    if text_only:
                        lines.append(str(item))
                    else:
                        # Indent multi-line content
                        for line in str(item).split('\n'):
                            lines.append(f"  {line}")
        
        if slide.get('speaker_notes'):
            if text_only:
                lines.append(f"[Notes: {slide['speaker_notes']}]")
            else:
                notes_preview = slide['speaker_notes']
                if len(notes_preview) > 200:
                    notes_preview = notes_preview[:200] + "..."
                lines.append(f"Speaker Notes: {notes_preview}")
        
        lines.append("")
    
    return "\n".join(lines)


def format_as_summary(presentation_info: dict) -> str:
    """Format presentation info as a brief summary."""
    lines = []
    
    lines.append(f"Presentation: {presentation_info['file_name']}")
    lines.append(f"Total Slides: {presentation_info['total_slides']}")
    
    summary = presentation_info.get('summary', {})
    if summary.get('slides_with_notes') is not None:
        lines.append(f"Slides with Notes: {summary['slides_with_notes']}")
    lines.append(f"Slides with Titles: {summary.get('slides_with_titles', 0)}")
    
    if presentation_info['metadata']['title']:
        lines.append(f"Document Title: {presentation_info['metadata']['title']}")
    
    lines.append("")
    lines.append("Slide Overview:")
    lines.append("-" * 40)
    
    for slide in presentation_info['slides']:
        title = slide['title'] or "(No title)"
        notes_indicator = " [has notes]" if slide.get('has_notes') else ""
        lines.append(f"  {slide['slide_number']:3}. {title}{notes_indicator}")
    
    return "\n".join(lines)


def main():
    parser = argparse.ArgumentParser(
        description="Read and extract information from PowerPoint presentations.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  %(prog)s presentation.pptx
  %(prog)s presentation.pptx --format json
  %(prog)s presentation.pptx --text-only
  %(prog)s presentation.pptx --slide 2
        """,
    )
    parser.add_argument(
        "input_file",
        help="Input PPTX file path",
    )
    parser.add_argument(
        "--format",
        choices=["text", "json", "summary"],
        default="summary",
        help="Output format (default: summary)",
    )
    parser.add_argument(
        "--text-only",
        action="store_true",
        help="Extract only text content (no metadata)",
    )
    parser.add_argument(
        "--slide",
        type=int,
        help="Get info for a specific slide (1-based index)",
    )
    parser.add_argument(
        "--include-notes",
        action="store_true",
        default=True,
        help="Include speaker notes in output (default: true)",
    )
    parser.add_argument(
        "--no-notes",
        action="store_true",
        help="Exclude speaker notes from output",
    )
    parser.add_argument(
        "--output",
        help="Write output to file instead of stdout",
    )
    
    args = parser.parse_args()
    
    # Check if file exists
    input_path = Path(args.input_file)
    if not input_path.exists():
        print(f"Error: File not found: {args.input_file}", file=sys.stderr)
        sys.exit(2)
    
    # Handle notes flag
    include_notes = not args.no_notes
    
    try:
        presentation_info = extract_presentation_content(args.input_file, include_notes)
    except Exception as e:
        print(f"Error: Failed to read presentation: {e}", file=sys.stderr)
        sys.exit(3)
    
    # Filter to specific slide if requested
    if args.slide:
        if args.slide < 1 or args.slide > len(presentation_info['slides']):
            print(f"Error: Slide {args.slide} out of range (1-{len(presentation_info['slides'])})", file=sys.stderr)
            sys.exit(1)
        
        # Keep only the requested slide
        presentation_info['slides'] = [presentation_info['slides'][args.slide - 1]]
    
    # Format output
    if args.text_only:
        output = format_as_text(presentation_info, text_only=True)
    elif args.format == "json":
        output = json.dumps(presentation_info, indent=2, ensure_ascii=False)
    elif args.format == "text":
        output = format_as_text(presentation_info)
    else:  # summary
        output = format_as_summary(presentation_info)
    
    # Write output
    if args.output:
        try:
            Path(args.output).write_text(output, encoding="utf-8")
            print(f"Output written to: {args.output}", file=sys.stderr)
        except Exception as e:
            print(f"Error: Failed to write output: {e}", file=sys.stderr)
            sys.exit(3)
    else:
        print(output)


if __name__ == "__main__":
    main()
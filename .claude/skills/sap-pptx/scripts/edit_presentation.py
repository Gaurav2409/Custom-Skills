#!/usr/bin/env python3
# /// script
# requires-python = ">=3.9"
# dependencies = [
#     "python-pptx>=0.6.21",
# ]
# ///
"""
Edit existing PowerPoint presentations.

Usage:
    uv run scripts/edit_presentation.py INPUT_FILE [OPTIONS]

Options:
    --output FILE         Output file path (default: overwrites input)
    --add-slide LAYOUT    Add a new slide with specified layout
    --slide NUMBER        Target slide for modifications (1-based)
    --title TEXT          Set/update slide title
    --body TEXT           Set/update slide body text
    --delete-slide NUM    Delete slide at position (1-based)
    --move-slide NUM      Move slide from position
    --to-position NUM     Move slide to position
    --add-image PATH      Add image to slide
    --image-position POS  Image position: left, center, right, full (default: center)
    --dry-run             Preview changes without saving
    --help                Show this help message

Exit Codes:
    0 - Success
    1 - Invalid arguments
    2 - File not found
    3 - Invalid layout name
    4 - Invalid slide number
    5 - Write error

Examples:
    uv run scripts/edit_presentation.py presentation.pptx --add-slide "Title and Text" --title "New Slide"
    uv run scripts/edit_presentation.py presentation.pptx --slide 2 --title "Updated Title"
    uv run scripts/edit_presentation.py presentation.pptx --delete-slide 3
    uv run scripts/edit_presentation.py presentation.pptx --move-slide 3 --to-position 1
"""

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


# Layout name to index mapping (0-based)
LAYOUT_MAP = {
    "Cover A": 0,
    "Cover B": 1,
    "Cover C": 2,
    "Cover D": 3,
    "Cover E": 4,
    "Cover F": 5,
    "Cover G": 6,
    "Cover H": 7,
    "Cover I": 8,
    "Cover J": 9,
    "Cover K": 10,
    "Cover L": 11,
    "Agenda A": 12,
    "Agenda B": 13,
    "Divider Page A": 14,
    "Divider Page B": 15,
    "Divider Page C": 16,
    "Divider Page D": 17,
    "Title Only": 18,
    "Title and Text": 19,
    "Title and Text: 2 Columns": 20,
    "Title and Text: 3 Columns": 21,
    "2 Columns - Text and Images": 22,
    "3 Columns - Text and Images": 23,
    "4 Columns - Text and Images": 24,
    "Title and Text with Image 1/3": 25,
    "Full Bleed Image": 26,
    "Text and Screenshot": 27,
    "Title and Content": 28,
    "Quote": 29,
    "Q & A": 30,
    "Thank You A": 31,
    "Thank You B": 32,
    "Blank": 33,
}


def get_layout_index(layout_name: str) -> int:
    """Get the layout index from layout name."""
    if layout_name in LAYOUT_MAP:
        return LAYOUT_MAP[layout_name]
    
    # Try case-insensitive match
    layout_lower = layout_name.lower()
    for name, index in LAYOUT_MAP.items():
        if name.lower() == layout_lower:
            return index
    
    return -1


def add_slide(prs: Presentation, layout_name: str, title: str | None = None,
              body: str | None = None) -> dict:
    """Add a new slide with the specified layout."""
    layout_index = get_layout_index(layout_name)
    if layout_index < 0:
        raise ValueError(f"Invalid layout name: {layout_name}")
    
    if layout_index >= len(prs.slide_layouts):
        raise ValueError(f"Layout index {layout_index} out of range.")
    
    slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)
    slide_number = len(prs.slides)
    
    # Set title if provided
    if title:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.is_placeholder:
                try:
                    ph_type = shape.placeholder_format.type
                    if ph_type in (1, 3):  # TITLE, CENTER_TITLE
                        shape.text = title
                        break
                except Exception:
                    pass
    
    # Set body if provided
    if body:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.is_placeholder:
                try:
                    ph_type = shape.placeholder_format.type
                    if ph_type in (2, 7):  # BODY, OBJECT
                        shape.text = body
                        break
                except Exception:
                    pass
    
    return {
        "action": "add_slide",
        "layout": layout_name,
        "slide_number": slide_number,
        "title": title,
    }


def update_slide(prs: Presentation, slide_number: int, title: str | None = None,
                 body: str | None = None) -> dict:
    """Update content on an existing slide."""
    if slide_number < 1 or slide_number > len(prs.slides):
        raise ValueError(f"Slide {slide_number} not found. Presentation has {len(prs.slides)} slides.")
    
    slide = prs.slides[slide_number - 1]
    changes = []
    
    # Update title if provided
    if title is not None:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.is_placeholder:
                try:
                    ph_type = shape.placeholder_format.type
                    if ph_type in (1, 3):  # TITLE, CENTER_TITLE
                        shape.text = title
                        changes.append("title")
                        break
                except Exception:
                    pass
    
    # Update body if provided
    if body is not None:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.is_placeholder:
                try:
                    ph_type = shape.placeholder_format.type
                    if ph_type in (2, 7):  # BODY, OBJECT
                        shape.text = body
                        changes.append("body")
                        break
                except Exception:
                    pass
    
    return {
        "action": "update_slide",
        "slide_number": slide_number,
        "changes": changes,
    }


def delete_slide(prs: Presentation, slide_number: int) -> dict:
    """Delete a slide at the specified position."""
    if slide_number < 1 or slide_number > len(prs.slides):
        raise ValueError(f"Slide {slide_number} not found. Presentation has {len(prs.slides)} slides.")
    
    rId = prs.slides._sldIdLst[slide_number - 1].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[slide_number - 1]
    
    return {
        "action": "delete_slide",
        "slide_number": slide_number,
        "remaining_slides": len(prs.slides),
    }


def move_slide(prs: Presentation, from_position: int, to_position: int) -> dict:
    """Move a slide from one position to another."""
    total_slides = len(prs.slides)
    
    if from_position < 1 or from_position > total_slides:
        raise ValueError(f"From position {from_position} invalid. Presentation has {total_slides} slides.")
    
    if to_position < 1 or to_position > total_slides:
        raise ValueError(f"To position {to_position} invalid. Must be between 1 and {total_slides}.")
    
    if from_position == to_position:
        return {
            "action": "move_slide",
            "from_position": from_position,
            "to_position": to_position,
            "note": "No change needed",
        }
    
    # Get the slide ID element
    slide_id = prs.slides._sldIdLst[from_position - 1]
    
    # Remove from current position
    prs.slides._sldIdLst.remove(slide_id)
    
    # Insert at new position
    prs.slides._sldIdLst.insert(to_position - 1, slide_id)
    
    return {
        "action": "move_slide",
        "from_position": from_position,
        "to_position": to_position,
    }


def add_image_to_slide(prs: Presentation, slide_number: int, image_path: str,
                       position: str = "center") -> dict:
    """Add an image to a slide."""
    if slide_number < 1 or slide_number > len(prs.slides):
        raise ValueError(f"Slide {slide_number} not found.")
    
    image_path_obj = Path(image_path)
    if not image_path_obj.exists():
        raise FileNotFoundError(f"Image not found: {image_path}")
    
    slide = prs.slides[slide_number - 1]
    
    # Calculate position based on slide size
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # Default image size (adjust as needed)
    img_width = Inches(5)
    
    if position == "left":
        left = Inches(0.5)
        top = Inches(2)
    elif position == "right":
        left = slide_width - img_width - Inches(0.5)
        top = Inches(2)
    elif position == "full":
        left = Inches(0)
        top = Inches(0)
        img_width = slide_width
    else:  # center
        left = (slide_width - img_width) / 2
        top = Inches(2)
    
    slide.shapes.add_picture(str(image_path), left, top, width=img_width)
    
    return {
        "action": "add_image",
        "slide_number": slide_number,
        "image_path": image_path,
        "position": position,
    }


def main():
    parser = argparse.ArgumentParser(
        description="Edit existing PowerPoint presentations.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  %(prog)s presentation.pptx --add-slide "Title and Text" --title "New Slide"
  %(prog)s presentation.pptx --slide 2 --title "Updated Title"
  %(prog)s presentation.pptx --delete-slide 3
  %(prog)s presentation.pptx --move-slide 3 --to-position 1
        """,
    )
    parser.add_argument(
        "input_file",
        help="Path to the PowerPoint file to edit",
    )
    parser.add_argument(
        "--output",
        help="Output file path (default: overwrites input)",
    )
    parser.add_argument(
        "--add-slide",
        metavar="LAYOUT",
        help="Add a new slide with specified layout",
    )
    parser.add_argument(
        "--slide",
        type=int,
        help="Target slide for modifications (1-based)",
    )
    parser.add_argument(
        "--title",
        help="Set/update slide title",
    )
    parser.add_argument(
        "--body",
        help="Set/update slide body text",
    )
    parser.add_argument(
        "--delete-slide",
        type=int,
        metavar="NUM",
        help="Delete slide at position (1-based)",
    )
    parser.add_argument(
        "--move-slide",
        type=int,
        metavar="NUM",
        help="Move slide from position",
    )
    parser.add_argument(
        "--to-position",
        type=int,
        metavar="NUM",
        help="Move slide to position",
    )
    parser.add_argument(
        "--add-image",
        metavar="PATH",
        help="Add image to slide",
    )
    parser.add_argument(
        "--image-position",
        choices=["left", "center", "right", "full"],
        default="center",
        help="Image position (default: center)",
    )
    parser.add_argument(
        "--dry-run",
        action="store_true",
        help="Preview changes without saving",
    )
    
    args = parser.parse_args()
    
    # Validate input file
    input_path = Path(args.input_file)
    if not input_path.exists():
        print(f"Error: File not found: {args.input_file}", file=sys.stderr)
        sys.exit(2)
    
    # Validate layout if adding slide
    if args.add_slide:
        if get_layout_index(args.add_slide) < 0:
            print(f"Error: Invalid layout name: {args.add_slide}", file=sys.stderr)
            print(f"       Use 'uv run scripts/list_layouts.py' to see available layouts.", file=sys.stderr)
            sys.exit(3)
    
    # Validate move arguments
    if args.move_slide and not args.to_position:
        print(f"Error: --move-slide requires --to-position", file=sys.stderr)
        sys.exit(1)
    
    try:
        prs = Presentation(str(input_path))
        results = []
        
        # Perform operations
        if args.delete_slide:
            result = delete_slide(prs, args.delete_slide)
            results.append(result)
        
        if args.move_slide and args.to_position:
            result = move_slide(prs, args.move_slide, args.to_position)
            results.append(result)
        
        if args.add_slide:
            result = add_slide(prs, args.add_slide, args.title, args.body)
            results.append(result)
        elif args.slide:
            # Update existing slide
            if args.title or args.body:
                result = update_slide(prs, args.slide, args.title, args.body)
                results.append(result)
            
            if args.add_image:
                result = add_image_to_slide(prs, args.slide, args.add_image, args.image_position)
                results.append(result)
        
        if not results:
            print("Warning: No operations specified.", file=sys.stderr)
            sys.exit(0)
        
        # Output results
        output_data = {
            "input_file": args.input_file,
            "operations": results,
            "dry_run": args.dry_run,
            "total_slides": len(prs.slides),
        }
        
        print(json.dumps(output_data, indent=2))
        
        # Save if not dry run
        if not args.dry_run:
            output_path = args.output or args.input_file
            prs.save(output_path)
            print(f"\nPresentation saved: {output_path}", file=sys.stderr)
        else:
            print(f"\nDry run - no changes saved.", file=sys.stderr)
        
    except ValueError as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(4)
    except FileNotFoundError as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(2)
    except Exception as e:
        print(f"Error: Failed to edit presentation: {e}", file=sys.stderr)
        sys.exit(5)


if __name__ == "__main__":
    main()